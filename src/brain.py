import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from .banana_gen import generate_response

class SyllabusBrain:
    def __init__(self, data_path='data/'):
        self.data_path = data_path
        if not os.path.exists(self.data_path): os.makedirs(self.data_path)
        print('Feeding AI... Please wait.')
        self.knowledge_base = self._load_data()

    def _load_data(self):
        content = ''
        for filename in os.listdir(self.data_path):
            file_path = os.path.join(self.data_path, filename)
            print(f'Reading: {filename}...')
            try:
                if filename.endswith('.txt'):
                    with open(file_path, 'r', encoding='utf-8') as f: content += f.read() + '\n'
                elif filename.endswith('.pdf'):
                    reader = PdfReader(file_path)
                    for page in reader.pages: content += page.extract_text() + '\n'
                elif filename.endswith('.docx'):
                    doc = Document(file_path)
                    content += '\n'.join([p.text for p in doc.paragraphs]) + '\n'
                elif filename.endswith('.pptx'):
                    prs = Presentation(file_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'): content += shape.text + '\n'
            except Exception as e: print(f'Error reading {filename}: {e}')
        print('Feeding complete!')
        return content

    def query(self, user_input):
        return generate_response(user_input, context=self.knowledge_base)
